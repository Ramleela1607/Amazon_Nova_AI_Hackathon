import io
import time
import json
import hashlib
import re
from typing import List, Dict, Any, Tuple, Optional

import streamlit as st
import pandas as pd
from pypdf import PdfReader
from PIL import Image

from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.lib import colors

from rag_index import RagIndex

from bedrock_utils import (
    ask_with_evidence,
    extract_fields_json,
    detect_doc_type,
    compare_docs,
    DOC_TYPES,
    recommend_rag_settings,
    nova_image_to_text,
    nova_image_insights_brief,
    generate_report_title,
    suggest_questions,
    generate_dashboard_insights_dynamic,  # AI dashboard (may fail JSON -> hybrid fallback)
)

# ------------------------------------------------------------
# OPTIONAL: Draggable tiles (UI only)
# ------------------------------------------------------------
try:
    from streamlit_elements import elements, dashboard as el_dashboard, mui
    ELEMENTS_OK = True
except Exception:
    ELEMENTS_OK = False


# ============================================================
# Page + Premium UI (bright readable + no top whitespace)
# ============================================================
st.set_page_config(page_title="Smart Document Copilot", layout="wide")

st.markdown(
    """
<style>
html, body { height: 100% !important; background: transparent !important; }
div[data-testid="stAppViewContainer"]{ background: transparent !important; }
header[data-testid="stHeader"]{ background: transparent !important; height: 0px !important; }
section.main > div{ padding-top: 0.6rem !important; }
.block-container{ padding-top: 0.7rem !important; }

:root{
  --ink: #07101f;
  --muted: rgba(7,16,31,0.70);
  --glass2: rgba(255,255,255,0.86);
  --stroke: rgba(15,23,42,0.12);
  --neonA: #4f46e5;
}

.stApp {
  background:
    radial-gradient(900px circle at 14% 16%, rgba(79,70,229,0.22), transparent 60%),
    radial-gradient(820px circle at 88% 14%, rgba(6,182,212,0.18), transparent 58%),
    radial-gradient(1000px circle at 70% 88%, rgba(34,197,94,0.14), transparent 62%),
    radial-gradient(900px circle at 22% 88%, rgba(249,115,22,0.10), transparent 62%),
    linear-gradient(135deg, rgba(255,255,255,0.94), rgba(245,247,255,0.86));
  background-attachment: fixed;
  color: var(--ink) !important;
  position: relative;
  overflow-x: hidden;
}
.stApp::before{
  content:"";
  position: fixed;
  inset: -20px;
  z-index: 0;
  pointer-events: none;
  background:
    radial-gradient(circle at 18% 32%, rgba(79,70,229,0.18), transparent 46%),
    radial-gradient(circle at 76% 22%, rgba(6,182,212,0.16), transparent 46%),
    radial-gradient(circle at 70% 82%, rgba(34,197,94,0.12), transparent 48%);
  filter: blur(18px);
  animation: floatBg 14s ease-in-out infinite alternate;
  opacity: 0.95;
}
@keyframes floatBg {
  0%   { transform: translate3d(-16px, -14px, 0) scale(1.02); }
  50%  { transform: translate3d(22px, 12px, 0) scale(1.06); }
  100% { transform: translate3d(-10px, 22px, 0) scale(1.03); }
}

.block-container,
section[data-testid="stSidebar"],
header, footer,
div[data-testid="stAppViewContainer"]{
  position: relative;
  z-index: 1;
}

html, body, [class*="css"], p, span, div, label {
  color: var(--ink) !important;
}
h1, h2, h3 {
  color: #061022 !important;
  text-shadow: 0 1px 0 rgba(255,255,255,0.65);
}

section[data-testid="stSidebar"]{
  background: rgba(255,255,255,0.82) !important;
  backdrop-filter: blur(14px);
  border-right: 1px solid var(--stroke);
}
section[data-testid="stSidebar"] *{ color: var(--ink) !important; }

div[data-testid="stExpander"] {
  border-radius: 18px;
  border: 1px solid var(--stroke);
  background: var(--glass2);
  box-shadow: 0 10px 28px rgba(2,6,23,0.08);
}

div[data-baseweb="input"] input, textarea {
  background: rgba(255,255,255,0.92) !important;
  color: var(--ink) !important;
  border-radius: 14px !important;
  border: 1px solid rgba(15,23,42,0.18) !important;
}

.stButton button {
  border-radius: 14px;
  border: 1px solid rgba(79,70,229,0.26);
  background: linear-gradient(135deg, rgba(79,70,229,0.10), rgba(6,182,212,0.08));
  color: var(--ink) !important;
  font-weight: 800;
  transition: transform 120ms ease, box-shadow 120ms ease;
}
.stButton button:hover {
  transform: translateY(-1px);
  box-shadow: 0 12px 28px rgba(79,70,229,0.16);
}

.pill {
  display:inline-block;
  padding:7px 11px;
  border-radius:999px;
  border:1px solid rgba(15,23,42,0.12);
  background: rgba(255,255,255,0.88);
  font-weight: 750;
  font-size: 13px;
}
.small-muted{ color: var(--muted) !important; font-size: 0.92rem; }

.tile {
  border-radius: 18px;
  border: 1px solid rgba(15,23,42,0.12);
  background: rgba(255,255,255,0.86);
  box-shadow: 0 12px 30px rgba(2,6,23,0.09);
  padding: 12px 14px;
}
.tile .title{ font-weight: 900; font-size: 0.95rem; margin-bottom: 6px; }
.tile .drag{
  cursor: move; user-select: none;
  display: inline-block;
  padding: 3px 8px;
  border-radius: 999px;
  border: 1px dashed rgba(79,70,229,0.30);
  background: rgba(79,70,229,0.06);
  margin-bottom: 10px;
}

a { color: #1d4ed8 !important; }
</style>
""",
    unsafe_allow_html=True,
)

st.title("📄 Smart Document Copilot")
st.markdown(
    "<div class='pill'><b>Amazon Nova</b> • Bedrock • Multimodal • RAG • Evidence • Dashboard • Compare • Report</div>"
    "<div class='small-muted'>Bright Neon Glass UI • Text readable • Optional draggable tiles</div>",
    unsafe_allow_html=True,
)

# ============================================================
# Helpers (Chunking, Extraction, OCR, Dates, Dashboard Mining)
# ============================================================

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 150) -> List[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        c = text[start:end].strip()
        if c:
            chunks.append(c)
        start += max(1, chunk_size - overlap)
    return chunks


def normalize_image_to_png_bytes(uploaded_img) -> bytes:
    img = Image.open(uploaded_img).convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def extract_text_from_pdf_basic(file_like) -> str:
    reader = PdfReader(file_like)
    texts = []
    for i, page in enumerate(reader.pages):
        t = page.extract_text() or ""
        t = t.strip()
        if t:
            texts.append(f"\n\n--- Page {i+1} ---\n{t}")
    return "".join(texts)


def pdf_pages_to_png_bytes(pdf_bytes: bytes, max_pages: int = 6) -> List[bytes]:
    try:
        import fitz  # PyMuPDF
    except Exception:
        return []

    pages = []
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    for i in range(min(max_pages, doc.page_count)):
        page = doc.load_page(i)
        pix = page.get_pixmap(dpi=180)
        pages.append(pix.tobytes("png"))
    return pages


def extract_text_from_pdf_with_ocr(uploaded_pdf, max_ocr_pages: int = 6) -> str:
    pdf_bytes = uploaded_pdf.getvalue()
    cache_key = "pdf_text:" + hashlib.md5(pdf_bytes[:250000]).hexdigest()

    if cache_key in st.session_state:
        return st.session_state[cache_key]

    basic_text = extract_text_from_pdf_basic(io.BytesIO(pdf_bytes))
    basic_len = len((basic_text or "").strip())
    combined = basic_text or ""

    if basic_len < 250:
        ocr_parts = []
        page_pngs = pdf_pages_to_png_bytes(pdf_bytes, max_pages=max_ocr_pages)
        if page_pngs:
            with st.spinner(f"🧠 OCR on PDF pages (first {len(page_pngs)} pages)..."):
                for idx, png_bytes in enumerate(page_pngs, start=1):
                    try:
                        t = nova_image_to_text(png_bytes, image_format="png")
                    except Exception:
                        t = ""
                    if t.strip():
                        ocr_parts.append(f"\n\n--- OCR Page {idx} ---\n{t.strip()}")
        if ocr_parts:
            combined = (combined + "\n\n" + "\n".join(ocr_parts)).strip()

    st.session_state[cache_key] = combined
    return combined


def extract_text_from_word(uploaded_docx) -> str:
    try:
        import docx
    except Exception:
        return ""
    try:
        b = uploaded_docx.getvalue()
        d = docx.Document(io.BytesIO(b))
        parts = []
        for p in d.paragraphs:
            if p.text.strip():
                parts.append(p.text.strip())
        for table in d.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        return "\n".join(parts).strip()
    except Exception:
        return ""


def extract_text_from_ppt(uploaded_pptx) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return ""
    try:
        b = uploaded_pptx.getvalue()
        prs = Presentation(io.BytesIO(b))
        parts = []
        for si, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                parts.append(f"--- Slide {si} ---\n" + "\n".join(slide_text))
        return "\n\n".join(parts).strip()
    except Exception:
        return ""


def read_excel_or_csv(uploaded_file) -> Tuple[str, List[Tuple[str, pd.DataFrame]]]:
    name = (uploaded_file.name or "").lower()
    raw = uploaded_file.getvalue()
    previews: List[Tuple[str, pd.DataFrame]] = []

    if name.endswith(".csv"):
        try:
            df = pd.read_csv(io.BytesIO(raw))
            previews.append(("CSV", df))
            text = df.to_csv(index=False)
            return text, previews
        except Exception:
            return "", []

    if name.endswith(".xlsx") or name.endswith(".xls"):
        try:
            xls = pd.ExcelFile(io.BytesIO(raw))
            for sheet in xls.sheet_names[:6]:
                df = xls.parse(sheet)
                previews.append((sheet, df))
            text_parts = []
            for sheet, df in previews[:4]:
                text_parts.append(f"--- SHEET: {sheet} ---\n{df.head(80).to_csv(index=False)}")
            return "\n\n".join(text_parts).strip(), previews
        except Exception:
            return "", []

    return "", []


def try_parse_number(value: str):
    if value is None:
        return None
    s = str(value).strip().replace(",", "")
    m = re.search(r"(-?\d+(?:\.\d+)?)", s)
    if not m:
        return None
    try:
        return float(m.group(1))
    except Exception:
        return None


def extract_dates_with_events(text: str, max_items: int = 140) -> List[Dict[str, str]]:
    if not text or not str(text).strip():
        return []

    t = str(text)
    t = t.replace("\r", "\n")
    t = t.replace("\u2013", "-").replace("\u2014", "-")
    t = re.sub(r"[ \t]+", " ", t)
    t = re.sub(r"\n{2,}", "\n", t)
    t = re.sub(r"(\d{1,2})\s*\n\s*([A-Za-z]{3,9})\s*\n\s*(\d{2,4})", r"\1 \2 \3", t)
    t = re.sub(r"\b(\d{1,2})(st|nd|rd|th)\b", r"\1", t, flags=re.IGNORECASE)

    months = r"(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Sept|Oct|Nov|Dec)[a-z]*"
    p_range = rf"\b({months}\s+\d{{4}})\s*-\s*(Present|{months}\s+\d{{4}})\b"
    p_iso = r"\b\d{4}[/-]\d{1,2}[/-]\d{1,2}\b"
    p_slash = r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b"
    p_mdy = rf"\b{months}\s+\d{{1,2}},?\s+\d{{2,4}}\b"
    p_dmy = rf"\b\d{{1,2}}\s+{months},?\s+\d{{2,4}}\b"
    p_my = rf"\b{months}\s+\d{{4}}\b"
    date_re = re.compile("|".join([p_range, p_iso, p_slash, p_mdy, p_dmy, p_my]), re.IGNORECASE)

    results: List[Dict[str, str]] = []
    seen = set()
    for line in t.split("\n"):
        ln = line.strip()
        if not ln:
            continue
        for m in date_re.finditer(ln):
            date_str = m.group(0).strip()
            event = ln[:m.start()].strip() or ln[m.end():].strip()
            event = re.sub(r"\s+", " ", event).strip(" -:•;|") or "Date mentioned"
            words = event.split()
            if len(words) > 16:
                event = " ".join(words[-16:])
            key = (event.lower(), date_str.lower())
            if key in seen:
                continue
            seen.add(key)
            results.append({"label": event, "value": date_str})
            if len(results) >= max_items:
                break
        if len(results) >= max_items:
            break
    return results


# ============================================================
# ✅ NEW (Excel-only): meaningful KPIs from DataFrames (NOT text)
# ============================================================
def local_dashboard_from_tables(previews: List[Tuple[str, pd.DataFrame]]) -> Dict[str, Any]:
    """
    Excel/CSV KPI builder (meaningful):
    - Rows/Cols/Missing%
    - Numeric columns: sum/avg/max (prefers financial-ish names)
    - Chart: top numeric sums or top categories
    """
    out = {"kpis": [], "charts": [], "table_preview": [], "derived_insights": []}
    if not previews:
        return out

    # Merge first few sheets for metrics (keep it small)
    frames = []
    for sheet, df in previews[:4]:
        if isinstance(df, pd.DataFrame) and not df.empty:
            x = df.copy()
            x.columns = [str(c) for c in x.columns]
            x["_sheet"] = str(sheet)
            frames.append(x)
    if not frames:
        return out

    big = pd.concat(frames, ignore_index=True, sort=False)

    rows = int(big.shape[0])
    cols = int(big.shape[1])
    miss = float(big.isna().mean().mean()) if rows and cols else 0.0

    out["kpis"].append({"label": "Rows", "value": f"{rows}", "note": ""})
    out["kpis"].append({"label": "Columns", "value": f"{cols}", "note": ""})
    out["kpis"].append({"label": "Missing Cells", "value": f"{miss*100:.1f}%", "note": ""})

    # Numeric columns
    numeric_cols = []
    for c in big.columns:
        if c == "_sheet":
            continue
        s = pd.to_numeric(big[c], errors="coerce")
        nn = int(s.notna().sum())
        if nn >= max(10, rows // 20):  # enough signal
            numeric_cols.append((c, s))

    def col_priority(name: str) -> int:
        n = name.lower()
        hits = [
            ("revenue", 5), ("sales", 5), ("amount", 4), ("total", 4),
            ("profit", 5), ("margin", 4), ("cost", 4), ("expense", 4),
            ("balance", 3), ("price", 3), ("qty", 3), ("quantity", 3),
        ]
        score = 1
        for k, w in hits:
            if k in n:
                score += w
        return score

    stats = []
    for c, s in numeric_cols:
        s2 = s.dropna()
        if s2.empty:
            continue
        stats.append(
            (
                col_priority(c),
                c,
                float(s2.sum()),
                float(s2.mean()),
                float(s2.max()),
                int(s2.count()),
            )
        )

    stats.sort(key=lambda x: (x[0], abs(x[2])), reverse=True)

    # Add best numeric KPIs (sum + avg)
    for pr, c, ssum, savg, smax, cnt in stats[:6]:
        out["kpis"].append({"label": f"{c} (Sum)", "value": f"{ssum:,.2f}", "note": f"n={cnt}"})
        out["kpis"].append({"label": f"{c} (Avg)", "value": f"{savg:,.2f}", "note": ""})

    # Derived insights
    out["derived_insights"].append(f"Detected {len(frames)} sheet(s), {rows:,} rows, {len(big.columns):,} columns.")
    if stats:
        out["derived_insights"].append(f"Top numeric column by sum: {stats[0][1]} ({stats[0][2]:,.2f}).")

    # Chart: top numeric sums (12 max)
    if stats:
        data = [{"x": c[:24], "y": float(ssum)} for _pr, c, ssum, _avg, _mx, _cnt in stats[:12]]
        out["charts"].append({"title": "Top Numeric Columns by Sum", "type": "bar", "data": data})

    # If a good category column exists, add chart for top categories
    cat_candidates = []
    for c in big.columns:
        if c == "_sheet":
            continue
        s = big[c].astype(str)
        uniq = s.nunique(dropna=True)
        if 2 <= uniq <= 30:
            cat_candidates.append((uniq, c))
    cat_candidates.sort(key=lambda x: x[0])
    if cat_candidates:
        cat_col = cat_candidates[0][1]
        vc = big[cat_col].astype(str).value_counts().head(10)
        out["charts"].append(
            {"title": f"Top Categories: {cat_col}", "type": "bar",
             "data": [{"x": str(k)[:24], "y": int(v)} for k, v in vc.items()]}
        )

    # Table preview: show first sheet head
    sname, df0 = previews[0]
    if isinstance(df0, pd.DataFrame) and not df0.empty:
        out["table_preview"] = df0.head(12).to_dict(orient="records")

    return out

def local_dashboard_from_text(text: str, max_items: int = 140) -> Dict[str, Any]:
    out = {"kpis": [], "charts": [], "table_preview": [], "derived_insights": []}
    if not text or not str(text).strip():
        return out

    t = str(text).replace("\r", "\n").replace("\u2013", "-").replace("\u2014", "-")
    t = re.sub(r"[ \t]+", " ", t)
    t = re.sub(r"\n{3,}", "\n\n", t)

    lines = [ln.strip() for ln in t.split("\n") if ln.strip()]

    joined = []
    i = 0
    while i < len(lines):
        ln = lines[i]
        if i + 1 < len(lines):
            nxt = lines[i + 1]
            if (not re.search(r"\d", ln)) and re.fullmatch(r"[₹$€]?\s*[\d,]+(\.\d+)?%?", nxt):
                joined.append(f"{ln}: {nxt}")
                i += 2
                continue
        joined.append(ln)
        i += 1

    months = r"(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Sept|Oct|Nov|Dec)[a-z]*"
    date_like = re.compile(rf"({months}\s+\d{{1,2}}|\d{{1,2}}[/-]\d{{1,2}}[/-]\d{{2,4}}|\b\d{{4}}\b)", re.IGNORECASE)

    def is_year(n: float) -> bool:
        return 1900 <= n <= 2100 and float(int(n)) == float(n)

    def bad_label(lbl: str) -> bool:
        s = (lbl or "").strip().lower()
        if not s or s in {"number", "value", "metric"}:
            return True
        if len(s) < 3:
            return True
        if any(x in s for x in ["date", "uat", "window", "timeline", "kickoff", "go-live", "present"]):
            return True
        return False

    candidates = []
    for ln in joined:
        m = re.search(r"^(.{2,60}?)\s*[:\-]\s*([₹$€]?\s*[\d,]+(?:\.\d+)?%?)\s*$", ln)
        if m:
            candidates.append((m.group(1).strip(), m.group(2).strip(), 3))

    for ln in joined:
        m = re.search(r"^(.{3,60}?)\s+([₹$€]?\s*[\d,]+(?:\.\d+)?%?)\s*$", ln)
        if m:
            if date_like.search(ln):
                continue
            candidates.append((m.group(1).strip(), m.group(2).strip(), 2))

    num_re = re.compile(r"(-?\d{1,3}(?:,\d{3})*(?:\.\d+)?|-?\d+(?:\.\d+)?)")
    for ln in joined:
        nums = list(num_re.finditer(ln))
        if not nums:
            continue
        if len(nums) >= 2 and len(ln) <= 220:
            out["table_preview"].append({"row": ln})
        for mm in nums[:2]:
            raw = mm.group(1)
            left = re.sub(r"\s+", " ", ln[:mm.start()].strip())
            label = " ".join(left.split()[-6:]) if left else "Number"
            candidates.append((label, raw, 1))

    numeric = []
    seen = set()
    for label, raw, prio in candidates[:max_items * 3]:
        num = try_parse_number(raw)
        if num is None:
            continue
        if is_year(num):
            continue
        if abs(num) < 10:
            continue
        if date_like.search(label):
            continue
        if bad_label(label):
            continue

        key = (label.lower().strip(), round(float(num), 6))
        if key in seen:
            continue
        seen.add(key)
        numeric.append((prio, label, float(num), raw))

    numeric_sorted = sorted(numeric, key=lambda x: (x[0], abs(x[2])), reverse=True)

    for prio, label, num, raw in numeric_sorted[:9]:
        out["kpis"].append({"label": label[:42], "value": raw, "note": ""})

    vals = [x[2] for x in numeric_sorted]
    if len(vals) >= 3:
        mn, mx = min(vals), max(vals)
        avg = sum(vals) / len(vals)
        out["derived_insights"].append(f"Detected {len(vals)} numeric metrics (filtered). Min={mn:g}, Max={mx:g}, Avg={avg:g}.")

    if numeric_sorted:
        data = [{"x": lab[:28], "y": float(v)} for _p, lab, v, _raw in numeric_sorted[:12]]
        out["charts"].append({"title": "Top Metrics (Filtered)", "type": "bar", "data": data})

    out["table_preview"] = out["table_preview"][:8]
    return out


def merge_ai_and_local(ai: Dict[str, Any], local: Dict[str, Any]) -> Dict[str, Any]:
    merged = dict(ai or {})
    if not isinstance(merged.get("kpis"), list) or not merged.get("kpis"):
        merged["kpis"] = local.get("kpis", [])
    if not isinstance(merged.get("charts"), list) or not merged.get("charts"):
        merged["charts"] = local.get("charts", [])
    if not isinstance(merged.get("table_preview"), list) or not merged.get("table_preview"):
        merged["table_preview"] = local.get("table_preview", [])
    if not isinstance(merged.get("derived_insights"), list) or not merged.get("derived_insights"):
        merged["derived_insights"] = local.get("derived_insights", [])

    if not merged.get("summary") or "could not be generated" in str(merged.get("summary", "")).lower():
        if local.get("kpis") or local.get("charts"):
            merged["summary"] = "Auto dashboard generated from detected signals in the document (hybrid)."
        else:
            merged["summary"] = "No strong numeric/table signals detected yet."

    merged.setdefault("doc_type_guess", "generic")
    merged.setdefault("risk_score", 0)
    merged.setdefault("risks", [])
    merged.setdefault("next_actions", [])
    merged.setdefault("table_preview", [])
    merged.setdefault("charts", [])
    merged.setdefault("kpis", [])
    merged.setdefault("derived_insights", [])
    return merged


def make_pdf_report(filename: str, title: str, sections: List[Tuple[str, str]], table_blocks: List[Tuple[str, pd.DataFrame]] = None) -> bytes:
    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Heading1"]), Spacer(1, 0.2 * inch)]

    for heading, body in sections:
        story.append(Paragraph(heading, styles["Heading2"]))
        story.append(Spacer(1, 0.08 * inch))
        safe = (body or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        story.append(Paragraph(safe.replace("\n", "<br/>"), styles["BodyText"]))
        story.append(Spacer(1, 0.14 * inch))

    if table_blocks:
        story.append(Spacer(1, 0.12 * inch))
        story.append(Paragraph("Appendix: Table Previews", styles["Heading2"]))
        story.append(Spacer(1, 0.08 * inch))
        for tname, df in table_blocks[:4]:
            story.append(Paragraph(str(tname), styles["Heading3"]))
            story.append(Spacer(1, 0.06 * inch))
            dfx = df.head(18).copy()
            data = [dfx.columns.tolist()] + dfx.astype(str).values.tolist()
            tb = Table(data, hAlign="LEFT")
            tb.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.black),
                        ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("FONTSIZE", (0, 0), (-1, -1), 8),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ]
                )
            )
            story.append(tb)
            story.append(Spacer(1, 0.14 * inch))

    path = f"/tmp/{filename}"
    SimpleDocTemplate(path).build(story)
    with open(path, "rb") as f:
        return f.read()


def reset_session():
    st.session_state["uploader_key"] = st.session_state.get("uploader_key", 0) + 1
    for k in [
        "single_rag", "index_fp",
        "last_doc_fp", "auto_rag_settings",
        "suggest_fp", "suggested_questions",
        "latest_answer", "latest_evidence", "latest_sources", "latest_q", "latest_rt",
        "rag_a", "rag_b",
        "pending_question",
        "active_upload_kind",
        "active_upload_fp",
    ]:
        st.session_state.pop(k, None)

    for kk in list(st.session_state.keys()):
        if str(kk).startswith(("dashboard:", "img_insights:", "dates:", "pdf_text:", "extracted_text:", "tables:")):
            st.session_state.pop(kk, None)

    st.session_state["_do_rerun"] = True


def build_index_if_needed(full_text: str, chunk_size: int, overlap: int):
    if not full_text or len(full_text.strip()) < 10:
        return

    fp_src = f"{full_text[:25000]}||cs={chunk_size}||ov={overlap}"
    new_fp = hashlib.md5(fp_src.encode("utf-8", errors="ignore")).hexdigest()

    if st.session_state.get("index_fp") == new_fp and isinstance(st.session_state.get("single_rag"), RagIndex):
        return

    chunks = chunk_text(full_text, chunk_size=chunk_size, overlap=overlap)
    with st.spinner("🚀 Auto-building index (Nova embeddings → FAISS)..."):
        rag = RagIndex(dim=1024)
        rag.add_chunks(chunks)

    st.session_state["single_rag"] = rag
    st.session_state["index_fp"] = new_fp


def run_single_question(user_q: str, top_k: int):
    rag = st.session_state.get("single_rag", None)
    if not isinstance(rag, RagIndex):
        st.error("Index not ready. Rebuild from sidebar or Reset.")
        st.stop()

    with st.spinner("Retrieving sources..."):
        hits, scores = rag.search(user_q, k=top_k)
        ctx = [h.text for h in hits]

    with st.spinner("Thinking with Nova Lite..."):
        t0 = time.time()
        ans = ask_with_evidence(user_q, ctx)
        rt = round(time.time() - t0, 2)

    st.session_state["latest_q"] = user_q
    st.session_state["latest_answer"] = (ans.get("answer") or "").strip()
    st.session_state["latest_evidence"] = (ans.get("evidence") or "").strip()
    st.session_state["latest_sources"] = list(zip(hits, scores))
    st.session_state["latest_rt"] = rt


# ============================================================
# Session init + Sidebar
# ============================================================
st.session_state.setdefault("uploader_key", 0)

st.sidebar.header("⚙️ Controls")
st.sidebar.button("🔄 Reset / New session", on_click=reset_session, use_container_width=True)

mode = st.sidebar.radio("Mode", ["Single Document", "Compare Two Documents"])
st.sidebar.markdown("---")

# ============================================================
# Mode: Single Document
# ============================================================
if mode == "Single Document":
    k = st.session_state.get("uploader_key", 0)

    st.session_state.setdefault("active_upload_kind", "")
    active_kind = st.session_state.get("active_upload_kind", "")

    def set_active(kind: str, fp: str):
        st.session_state["active_upload_kind"] = kind
        st.session_state["active_upload_fp"] = fp

    def clear_active_if_none(pdf_u, img_u, xl_u, doc_u, ppt_u):
        if pdf_u is None and img_u is None and xl_u is None and doc_u is None and ppt_u is None:
            st.session_state["active_upload_kind"] = ""
            st.session_state.pop("active_upload_fp", None)

    st.markdown("### 1) Upload any document")
    cA, cB = st.columns([2, 1], vertical_alignment="top")

    with cA:
        uploaded_pdf = st.file_uploader(
            "📄 Upload PDF", type=["pdf"], key=f"pdf_uploader_{k}",
            disabled=bool(active_kind and active_kind != "pdf"),
        )
        uploaded_img = st.file_uploader(
            "🖼️ Upload Image (PNG/JPG/WebP)", type=["png", "jpg", "jpeg", "webp"],
            key=f"img_uploader_{k}",
            disabled=bool(active_kind and active_kind != "img"),
        )
        uploaded_xl = st.file_uploader(
            "📊 Upload Excel/CSV", type=["xlsx", "xls", "csv"],
            key=f"xl_uploader_{k}",
            disabled=bool(active_kind and active_kind != "xl"),
        )
        uploaded_doc = st.file_uploader(
            "📝 Upload Word (DOCX)", type=["docx"],
            key=f"doc_uploader_{k}",
            disabled=bool(active_kind and active_kind != "docx"),
        )
        uploaded_ppt = st.file_uploader(
            "📽️ Upload PPT (PPTX)", type=["pptx"],
            key=f"ppt_uploader_{k}",
            disabled=bool(active_kind and active_kind != "pptx"),
        )
        user_text = st.text_area("✍️ Paste extra notes (optional)", height=90)

    with cB:
        st.markdown("#### Tips")
        st.info(
            "• Upload **one** file at a time (others auto-disable)\n"
            "• Scanned PDFs: OCR runs automatically\n"
            "• Excel: shows table previews + meaningful KPIs"
        )

    clear_active_if_none(uploaded_pdf, uploaded_img, uploaded_xl, uploaded_doc, uploaded_ppt)

    if uploaded_pdf is None and uploaded_img is None and uploaded_xl is None and uploaded_doc is None and uploaded_ppt is None and not user_text.strip():
        st.info("Upload a PDF/Image/Excel/Word/PPT (or paste text) → Dashboard + Questions + Chat.")
        st.stop()

    extracted_text_parts: List[str] = []
    table_previews: List[Tuple[str, pd.DataFrame]] = []
    insights_block = ""
    upload_fp = ""

    if uploaded_pdf is not None:
        b = uploaded_pdf.getvalue()
        upload_fp = "pdf:" + hashlib.md5(b[:250000]).hexdigest()
        set_active("pdf", upload_fp)

        with st.spinner("Extracting PDF text/OCR..."):
            pdf_text = extract_text_from_pdf_with_ocr(uploaded_pdf, max_ocr_pages=6)
        if pdf_text.strip():
            extracted_text_parts.append("=== PDF TEXT ===\n" + pdf_text)

    if uploaded_img is not None:
        img_bytes = normalize_image_to_png_bytes(uploaded_img)
        upload_fp = "img:" + hashlib.md5(img_bytes[:200000]).hexdigest()
        set_active("img", upload_fp)

        img_cache_key = f"img_insights:{upload_fp}"
        if img_cache_key not in st.session_state:
            with st.spinner("Generating quick image insights..."):
                try:
                    st.session_state[img_cache_key] = nova_image_insights_brief(img_bytes, image_format="png") or ""
                except Exception:
                    st.session_state[img_cache_key] = ""

        insights_block = st.session_state.get(img_cache_key, "") or ""
        try:
            with st.spinner("OCR on image..."):
                ocr_text = nova_image_to_text(img_bytes, image_format="png")
        except Exception:
            ocr_text = ""

        if ocr_text.strip():
            extracted_text_parts.append("=== IMAGE TEXT ===\n" + ocr_text)
        if insights_block.strip():
            extracted_text_parts.append("=== IMAGE INSIGHTS ===\n" + insights_block)

    if uploaded_xl is not None:
        raw = uploaded_xl.getvalue()
        upload_fp = "xl:" + hashlib.md5(raw[:250000]).hexdigest()
        set_active("xl", upload_fp)

        with st.spinner("Reading Excel/CSV..."):
            xl_text, previews = read_excel_or_csv(uploaded_xl)
        if xl_text.strip():
            extracted_text_parts.append("=== EXCEL/CSV TEXT ===\n" + xl_text)
        table_previews = previews

    if uploaded_doc is not None:
        raw = uploaded_doc.getvalue()
        upload_fp = "docx:" + hashlib.md5(raw[:250000]).hexdigest()
        set_active("docx", upload_fp)

        with st.spinner("Reading Word document..."):
            doc_text = extract_text_from_word(uploaded_doc)
        if doc_text.strip():
            extracted_text_parts.append("=== WORD TEXT ===\n" + doc_text)

    if uploaded_ppt is not None:
        raw = uploaded_ppt.getvalue()
        upload_fp = "pptx:" + hashlib.md5(raw[:250000]).hexdigest()
        set_active("pptx", upload_fp)

        with st.spinner("Reading PowerPoint..."):
            ppt_text = extract_text_from_ppt(uploaded_ppt)
        if ppt_text.strip():
            extracted_text_parts.append("=== PPT TEXT ===\n" + ppt_text)

    if user_text.strip():
        extracted_text_parts.append("=== USER NOTES ===\n" + user_text.strip())

    full_text = "\n\n".join(extracted_text_parts).strip()

    with st.expander("🔎 Debug + Downloads", expanded=False):
        st.write("Extracted text length:", len(full_text))
        st.download_button(
            "⬇️ Download extracted text (.txt)",
            data=full_text.encode("utf-8", errors="ignore"),
            file_name="extracted_text.txt",
            mime="text/plain",
            use_container_width=True,
        )
        st.caption("Preview:")
        st.code((full_text[:1200] + "...") if len(full_text) > 1200 else full_text)

    if insights_block.strip():
        st.markdown("### 💡 Image Insights")
        st.markdown(insights_block.replace("\n", "  \n"))

    if table_previews:
        st.markdown("### 📊 Excel Preview Tables")
        for sheet_name, df in table_previews[:4]:
            st.markdown(f"**{sheet_name}**")
            st.dataframe(df.head(200), use_container_width=True)

    doc_fp = hashlib.md5((upload_fp + "||" + full_text[:25000]).encode("utf-8", errors="ignore")).hexdigest()

    dates_key = f"dates:{doc_fp}"
    if dates_key not in st.session_state:
        st.session_state[dates_key] = extract_dates_with_events(full_text, max_items=140)
    local_dates = st.session_state.get(dates_key, [])

    # =========================
    # Executive Dashboard
    # =========================
    st.subheader("📊 Executive Dashboard")
    dash_key = f"dashboard:{doc_fp}"

    dcol1, dcol2, dcol3 = st.columns([1, 1, 3])
    with dcol1:
        if st.button("🔄 Refresh Dashboard", use_container_width=True):
            st.session_state.pop(dash_key, None)
            st.rerun()
    with dcol2:
        if st.button("📄 Download Dashboard Report (PDF)", use_container_width=True):
            st.session_state["download_report_now"] = True

    # Local deterministic mining
    local_dash = local_dashboard_from_text(full_text)

    # ✅ Excel-only: replace local KPIs with DataFrame-based KPIs
    if st.session_state.get("active_upload_kind") == "xl" and table_previews:
        table_dash = local_dashboard_from_tables(table_previews)
        # prefer table-derived signals for Excel
        local_dash = merge_ai_and_local(table_dash, local_dash)

    if dash_key not in st.session_state:
        with st.spinner("Analyzing with Nova + Hybrid signals..."):
            ai_dash = {}
            if len(full_text.strip()) >= 80:
                try:
                    ai_dash = generate_dashboard_insights_dynamic(full_text) or {}
                except Exception:
                    ai_dash = {
                        "doc_type_guess": "generic",
                        "summary": "Dashboard AI failed; using hybrid local signals.",
                        "kpis": [], "derived_insights": [], "charts": [],
                        "table_preview": [], "risk_score": 0, "risks": [], "next_actions": [],
                    }
            else:
                ai_dash = {
                    "doc_type_guess": "generic",
                    "summary": "Not enough extracted text. Try clearer scan or OCR-friendly upload.",
                    "kpis": [], "derived_insights": [], "charts": [],
                    "table_preview": [], "risk_score": 0, "risks": [], "next_actions": [],
                }

            st.session_state[dash_key] = merge_ai_and_local(ai_dash, local_dash)

    dashboard = st.session_state.get(dash_key, {}) or {}
    summary = dashboard.get("summary", "")
    doc_type_guess = dashboard.get("doc_type_guess", "generic")
    risk_score = int(dashboard.get("risk_score", 0) or 0)
    kpis = dashboard.get("kpis", []) or []
    derived = dashboard.get("derived_insights", []) or []
    charts = dashboard.get("charts", []) or []
    dash_table_preview = dashboard.get("table_preview", []) or []
    risks = dashboard.get("risks", []) or []
    next_actions = dashboard.get("next_actions", []) or []

    # Draggable tiles (UI only)
    if ELEMENTS_OK:
        st.markdown("<div class='small-muted'>Drag tiles by the <b>Drag</b> handle to rearrange.</div>", unsafe_allow_html=True)
        st.session_state.setdefault(
            "dash_layout",
            [
                el_dashboard.Item("tile_doc", 0, 0, 4, 2),
                el_dashboard.Item("tile_risk", 4, 0, 4, 2),
                el_dashboard.Item("tile_sum", 8, 0, 4, 2),
                el_dashboard.Item("tile_kpi", 0, 2, 6, 5),
                el_dashboard.Item("tile_chart", 6, 2, 6, 5),
                el_dashboard.Item("tile_dates", 0, 7, 7, 5),
                el_dashboard.Item("tile_table", 7, 7, 5, 5),
            ],
        )

        with elements("dash"):
            with el_dashboard.Grid(st.session_state["dash_layout"], draggableHandle=".drag"):
                with mui.Paper(key="tile_doc", className="tile", elevation=0):
                    mui.Typography("Document Type", className="title")
                    mui.Typography("Drag", className="drag")
                    mui.Typography(str(doc_type_guess), variant="h5")

                with mui.Paper(key="tile_risk", className="tile", elevation=0):
                    mui.Typography("Risk Score", className="title")
                    mui.Typography("Drag", className="drag")
                    mui.Typography(f"{risk_score}/100", variant="h5")

                with mui.Paper(key="tile_sum", className="tile", elevation=0):
                    mui.Typography("Executive Summary", className="title")
                    mui.Typography("Drag", className="drag")
                    mui.Typography(summary or "—", variant="body2")

                with mui.Paper(key="tile_kpi", className="tile", elevation=0):
                    mui.Typography("KPIs", className="title")
                    mui.Typography("Drag", className="drag")
                    if kpis:
                        for x in kpis[:12]:
                            mui.Typography(f"• {x.get('label','KPI')}: {x.get('value','-')}", variant="body2")
                    else:
                        mui.Typography("No KPIs detected.", variant="body2")

                    if derived:
                        mui.Divider(sx={"my": 1})
                        mui.Typography("Derived Insights", fontWeight=900, variant="body2")
                        for d in derived[:6]:
                            mui.Typography(f"• {d}", variant="body2")

                with mui.Paper(key="tile_chart", className="tile", elevation=0):
                    mui.Typography("Charts", className="title")
                    mui.Typography("Drag", className="drag")
                    mui.Typography("Charts render below (Streamlit) for stability.", variant="body2")

                with mui.Paper(key="tile_dates", className="tile", elevation=0):
                    mui.Typography("Key Dates", className="title")
                    mui.Typography("Drag", className="drag")
                    if local_dates:
                        for d in local_dates[:10]:
                            mui.Typography(f"• {d.get('value','')} — {d.get('label','')}", variant="body2")
                    else:
                        mui.Typography("No dates detected.", variant="body2")

                with mui.Paper(key="tile_table", className="tile", elevation=0):
                    mui.Typography("Table Preview", className="title")
                    mui.Typography("Drag", className="drag")
                    if table_previews:
                        mui.Typography("Excel previews shown above/below.", variant="body2")
                    else:
                        mui.Typography("Preview shown below.", variant="body2")

    else:
        c1, c2, c3 = st.columns(3)
        c1.metric("Doc Type (Nova)", str(doc_type_guess))
        c2.metric("Risk Score", f"{risk_score}/100")
        with c3:
            st.caption("Risk Meter")
            st.progress(min(max(risk_score, 0), 100))

        if summary.strip():
            st.markdown("### 🧾 Executive Summary")
            st.markdown(summary)

        if kpis:
            st.markdown("### 🔑 KPIs")
            cols = st.columns(3)
            for i, kpi in enumerate(kpis[:9]):
                with cols[i % 3]:
                    st.metric(
                        str(kpi.get("label", "KPI"))[:40],
                        str(kpi.get("value", "-"))[:28],
                        (str(kpi.get("note", ""))[:38] if kpi.get("note") else None),
                    )

        if derived:
            st.markdown("### ✨ Derived Insights")
            for d in derived[:10]:
                st.markdown(f"- {d}")

    if charts:
        st.markdown("### 📈 Auto Charts")
        for ch in charts[:4]:
            title = ch.get("title", "Chart")
            ctype = ch.get("type", "bar")
            data = ch.get("data", []) or []
            st.markdown(f"**{title}**")
            if not data:
                st.caption("No chart data available.")
                continue

            dfc = pd.DataFrame(data)
            if "x" in dfc.columns and "y" in dfc.columns:
                dfc["x"] = dfc["x"].astype(str)
                dfc["y"] = pd.to_numeric(dfc["y"], errors="coerce").fillna(0)
                if ctype == "line":
                    st.line_chart(dfc.set_index("x")[["y"]])
                else:
                    st.bar_chart(dfc.set_index("x")[["y"]])
            else:
                st.dataframe(dfc, use_container_width=True)

    if table_previews:
        pass
    elif dash_table_preview:
        st.markdown("### 🧩 Table Preview")
        st.dataframe(pd.DataFrame(dash_table_preview), use_container_width=True)

    st.markdown("### 🗓️ Key Dates Timeline (AI Structured)")
    if local_dates:
        df_dates = pd.DataFrame(local_dates).rename(columns={"label": "Event", "value": "Date"})
        st.dataframe(df_dates, use_container_width=True)
    else:
        st.caption("No dates detected (often means extracted text is empty or dates are embedded in images without OCR).")

    if risks:
        st.markdown("### ⚠️ Risks")
        for r in risks[:8]:
            st.markdown(f"- {r}")

    if next_actions:
        st.markdown("### ✅ Next Actions")
        for a in next_actions[:8]:
            st.markdown(f"- {a}")

    st.divider()

    if st.session_state.pop("download_report_now", False):
        with st.spinner("Generating dashboard report PDF..."):
            title = "Smart Document Dashboard Report"
            sections = [
                ("Document Type", str(doc_type_guess)),
                ("Executive Summary", summary or "-"),
                ("Risk Score", f"{risk_score}/100"),
                ("Derived Insights", "\n".join(derived[:12]) if derived else "-"),
                ("KPIs", "\n".join([f"- {x.get('label','KPI')}: {x.get('value','-')}" for x in kpis[:12]]) if kpis else "-"),
                ("Key Dates", "\n".join([f"- {d.get('label','Event')} → {d.get('value','Date')}" for d in local_dates[:40]]) if local_dates else "-"),
                ("Latest Q&A", (
                    f"Question: {st.session_state.get('latest_q','-')}\n\n"
                    f"Answer:\n{st.session_state.get('latest_answer','-')}\n\n"
                    f"Evidence:\n{st.session_state.get('latest_evidence','-')}"
                ) if st.session_state.get("latest_answer") else "No Q&A yet."),
            ]
            pdf_bytes = make_pdf_report(
                "dashboard_report.pdf",
                title,
                sections,
                table_blocks=(table_previews if table_previews else []),
            )
        st.download_button(
            "⬇️ Download PDF Report",
            data=pdf_bytes,
            file_name="Smart_Document_Dashboard_Report.pdf",
            mime="application/pdf",
            use_container_width=True,
        )

    if st.session_state.get("last_doc_fp") != doc_fp:
        st.session_state["last_doc_fp"] = doc_fp
        with st.spinner("⚙️ Nova is auto-optimizing retrieval settings..."):
            st.session_state["auto_rag_settings"] = recommend_rag_settings(full_text)
        st.session_state.pop("suggest_fp", None)
        st.session_state.pop("suggested_questions", None)
        st.session_state.pop("index_fp", None)

    rec = st.session_state.get("auto_rag_settings", {"chunk_size": 1000, "overlap": 150, "top_k": 4})
    auto_chunk_size, auto_overlap, auto_top_k = rec["chunk_size"], rec["overlap"], rec["top_k"]

    st.sidebar.subheader("🔎 Retrieval settings")
    use_auto = st.sidebar.toggle("Use Nova auto-optimized settings", value=True)
    if use_auto:
        chunk_size, overlap, top_k = auto_chunk_size, auto_overlap, auto_top_k
        st.sidebar.caption(f"Auto: chunk={chunk_size}, overlap={overlap}, top_k={top_k}")
    else:
        chunk_size = st.sidebar.slider("Chunk size (chars)", 300, 2000, int(auto_chunk_size), 50)
        overlap = st.sidebar.slider("Overlap (chars)", 0, 400, int(auto_overlap), 25)
        top_k = st.sidebar.slider("Top-K sources", 2, 8, int(auto_top_k), 1)

    if st.sidebar.button("♻️ Rebuild index now", use_container_width=True):
        st.session_state.pop("index_fp", None)

    build_index_if_needed(full_text, chunk_size=chunk_size, overlap=overlap)

    st.divider()

    st.subheader("2) Extract key fields (JSON)")
    doc_type = st.selectbox("Document type", DOC_TYPES, index=0)
    if st.button("🧾 Extract key fields as JSON", use_container_width=True):
        with st.spinner("Extracting..."):
            out = extract_fields_json(full_text, doc_type=doc_type)
        st.code(out, language="json")

    st.divider()

    st.markdown("### ✨ Nova-suggested questions (auto from your document)")
    suggest_fp = f"{doc_fp}:General"
    if st.session_state.get("suggest_fp") != suggest_fp:
        st.session_state["suggest_fp"] = suggest_fp
        with st.spinner("Generating suggested questions..."):
            st.session_state["suggested_questions"] = suggest_questions(full_text, user_interest="General", n=6)

    qs = st.session_state.get("suggested_questions", []) or []
    if qs:
        cols = st.columns(3)
        for i, q in enumerate(qs):
            with cols[i % 3]:
                if st.button(q, use_container_width=True, key=f"dynq_{doc_fp}_{i}"):
                    st.session_state["pending_question"] = q
                    st.rerun()
    else:
        st.caption("No suggestions generated (try a clearer upload / more text).")

    if st.button("🔄 Refresh questions", use_container_width=True):
        with st.spinner("Refreshing questions..."):
            st.session_state["suggested_questions"] = suggest_questions(full_text, user_interest="General", n=6)
        st.rerun()

    st.divider()

    st.subheader("3) Chat with your document")
    st.markdown("#### 💬 Ask a question")

    pending = st.session_state.pop("pending_question", "")
    if pending:
        run_single_question(pending, top_k=top_k)

    with st.form("ask_form", clear_on_submit=True):
        q_text = st.text_input(
            "Type your question",
            placeholder="e.g., What are the key KPIs and what do they indicate?",
            label_visibility="collapsed",
        )
        ask_clicked = st.form_submit_button("Ask", use_container_width=True)

    if ask_clicked and q_text.strip():
        run_single_question(q_text.strip(), top_k=top_k)

    if st.session_state.get("latest_answer"):
        st.markdown("### ✅ Latest Answer")
        st.markdown(f"**Question:** {st.session_state.get('latest_q','')}")
        st.markdown(st.session_state["latest_answer"] or "I don't know based on the document.")
        ev = st.session_state.get("latest_evidence", "")
        if ev:
            with st.expander("📌 Evidence"):
                st.markdown(ev)
        st.caption(f"⏱ Response time: {st.session_state.get('latest_rt', 0)} sec")

else:
    st.subheader("🆚 Compare Two PDFs")
    k = st.session_state.get("uploader_key", 0)

    c1, c2 = st.columns(2)
    with c1:
        up_a = st.file_uploader("Upload PDF (Doc A)", type=["pdf"], key=f"docA_{k}")
    with c2:
        up_b = st.file_uploader("Upload PDF (Doc B)", type=["pdf"], key=f"docB_{k}")

    if up_a is None or up_b is None:
        st.info("Upload both PDFs → Ask a comparison question.")
        st.stop()

    text_a = extract_text_from_pdf_with_ocr(up_a, max_ocr_pages=4)
    text_b = extract_text_from_pdf_with_ocr(up_b, max_ocr_pages=4)

    with st.spinner("⚙️ Auto-tuning retrieval settings for comparison..."):
        rec = recommend_rag_settings(text_a + "\n\n" + text_b)

    chunk_size, overlap, top_k = rec["chunk_size"], rec["overlap"], rec["top_k"]

    chunks_a = chunk_text(text_a, chunk_size=chunk_size, overlap=overlap)
    chunks_b = chunk_text(text_b, chunk_size=chunk_size, overlap=overlap)

    if "rag_a" not in st.session_state or "rag_b" not in st.session_state:
        with st.spinner("🚀 Auto-building indexes for Doc A & Doc B..."):
            rag_a = RagIndex(dim=1024)
            rag_a.add_chunks(chunks_a)
            rag_b = RagIndex(dim=1024)
            rag_b.add_chunks(chunks_b)
        st.session_state["rag_a"] = rag_a
        st.session_state["rag_b"] = rag_b

    q = st.text_input("Comparison question", placeholder="e.g., Which doc has stronger KPIs and why?")
    if st.button("🆚 Compare", use_container_width=True):
        if not q.strip():
            st.error("Type a comparison question.")
            st.stop()

        with st.spinner("Retrieving from Doc A..."):
            hits_a, _ = st.session_state.rag_a.search(q, k=top_k)
            ctx_a = [h.text for h in hits_a]

        with st.spinner("Retrieving from Doc B..."):
            hits_b, _ = st.session_state.rag_b.search(q, k=top_k)
            ctx_b = [h.text for h in hits_b]

        with st.spinner("Comparing with Nova Lite..."):
            out = compare_docs(q, ctx_a, ctx_b, label_a="Doc A", label_b="Doc B")

        st.markdown("### ✅ Comparison result")
        st.write(out)

if st.session_state.pop("_do_rerun", False):
    st.rerun()

